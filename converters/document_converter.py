import os
import asyncio
from PIL import Image
import img2pdf
from PyPDF2 import PdfReader, PdfWriter
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter, A4
import tempfile
import pandas as pd
from config import Config
import subprocess
import sys
import fitz  # PyMuPDF for better PDF handling
import logging
from ebooklib import epub
import html

logger = logging.getLogger(__name__)

class DocumentConverter:
    def __init__(self):
        self.supported_conversions = {
            'pdf_to_docx': (['pdf'], 'docx'),
            'pdf_to_images': (['pdf'], 'images'),
            'docx_to_pdf': (['docx', 'doc'], 'pdf'),
            'images_to_pdf': (['jpg', 'jpeg', 'png', 'webp', 'bmp', 'tiff'], 'pdf'),
            'excel_to_pdf': (['xlsx', 'xls', 'csv'], 'pdf'),
            'ppt_to_pdf': (['pptx', 'ppt'], 'pdf'),
            'html_to_pdf': (['html', 'htm'], 'pdf'),
            'txt_to_pdf': (['txt'], 'pdf'),
            'markdown_to_pdf': (['md', 'markdown'], 'pdf'),
            'pdf_to_txt': (['pdf'], 'txt'),
            'pdf_to_html': (['pdf'], 'html'),
            'pdf_to_epub': (['pdf'], 'epub'),
        }
    
    async def convert_document(self, input_path, output_format):
        """Main document conversion method that routes to specific converters"""
        try:
            input_extension = input_path.split('.')[-1].lower()
            
            logger.info(f"Converting document: {input_extension} -> {output_format}")
            
            if input_extension == 'pdf':
                if output_format == 'txt':
                    output_path = input_path.rsplit('.', 1)[0] + '.txt'
                    return await self.convert_pdf_to_txt(input_path, output_path)
                elif output_format == 'docx':
                    output_path = input_path.rsplit('.', 1)[0] + '.docx'
                    return await self.convert_pdf_to_docx(input_path, output_path)
                elif output_format in ['jpg', 'jpeg', 'png', 'webp']:
                    return await self.convert_pdf_to_image(input_path, output_format)
                elif output_format == 'html':
                    output_path = input_path.rsplit('.', 1)[0] + '.html'
                    return await self.convert_pdf_to_html(input_path, output_path)
                elif output_format == 'epub':
                    output_path = input_path.rsplit('.', 1)[0] + '.epub'
                    return await self.convert_pdf_to_epub(input_path, output_path)
                else:
                    raise Exception(f"PDF to {output_format} conversion not supported")
            
            elif input_extension in ['docx', 'doc']:
                if output_format == 'pdf':
                    output_path = input_path.rsplit('.', 1)[0] + '.pdf'
                    return await self.convert_docx_to_pdf(input_path, output_path)
                elif output_format == 'txt':
                    output_path = input_path.rsplit('.', 1)[0] + '.txt'
                    return await self.convert_docx_to_txt(input_path, output_path)
                else:
                    raise Exception(f"DOCX to {output_format} conversion not supported")
            
            elif input_extension == 'txt':
                if output_format == 'pdf':
                    output_path = input_path.rsplit('.', 1)[0] + '.pdf'
                    return await self.convert_txt_to_pdf(input_path, output_path)
                else:
                    raise Exception(f"TXT to {output_format} conversion not supported")
            
            elif input_extension in ['xlsx', 'xls', 'csv']:
                if output_format == 'pdf':
                    output_path = input_path.rsplit('.', 1)[0] + '.pdf'
                    return await self.convert_excel_to_pdf(input_path, output_path)
                else:
                    raise Exception(f"Excel to {output_format} conversion not supported")
            
            elif input_extension in ['pptx', 'ppt']:
                if output_format == 'pdf':
                    output_path = input_path.rsplit('.', 1)[0] + '.pdf'
                    return await self.convert_ppt_to_pdf(input_path, output_path)
                else:
                    raise Exception(f"PowerPoint to {output_format} conversion not supported")
            
            elif input_extension in ['html', 'htm']:
                if output_format == 'pdf':
                    output_path = input_path.rsplit('.', 1)[0] + '.pdf'
                    return await self.convert_html_to_pdf(input_path, output_path)
                else:
                    raise Exception(f"HTML to {output_format} conversion not supported")
            
            elif input_extension in ['md', 'markdown']:
                if output_format == 'pdf':
                    output_path = input_path.rsplit('.', 1)[0] + '.pdf'
                    return await self.convert_markdown_to_pdf(input_path, output_path)
                else:
                    raise Exception(f"Markdown to {output_format} conversion not supported")
            
            else:
                raise Exception(f"Document conversion from {input_extension} to {output_format} not supported")
                
        except Exception as e:
            logger.error(f"Document conversion error: {e}")
            raise
    
    async def convert_pdf_to_txt(self, input_path, output_path):
        """Convert PDF to TXT using PyMuPDF"""
        try:
            doc = fitz.open(input_path)
            text_content = ""
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text = page.get_text()
                text_content += f"--- Page {page_num + 1} ---\n{text}\n\n"
            
            doc.close()
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text_content)
            
            return output_path
        except Exception as e:
            raise Exception(f"PDF to TXT conversion failed: {str(e)}")
    
    async def convert_pdf_to_docx(self, input_path, output_path):
        """Convert PDF to DOCX using pdf2docx"""
        try:
            from pdf2docx import Converter
            
            cv = Converter(input_path)
            cv.convert(output_path)
            cv.close()
            
            return output_path
        except Exception as e:
            # Fallback to text extraction
            return await self._fallback_pdf_to_docx(input_path, output_path)
    
    async def convert_pdf_to_image(self, input_path, output_format):
        """Convert PDF to image"""
        try:
            from pdf2image import convert_from_path
            
            images = convert_from_path(input_path, dpi=150)
            if images:
                output_path = input_path.rsplit('.', 1)[0] + f'.{output_format}'
                # Use correct format name for PIL
                pil_format = 'JPEG' if output_format in ['jpg', 'jpeg'] else output_format.upper()
                images[0].save(output_path, format=pil_format, quality=95)
                return output_path
            else:
                raise Exception("No pages found in PDF")
        except Exception as e:
            raise Exception(f"PDF to image conversion failed: {str(e)}")
    
    async def convert_pdf_to_html(self, input_path, output_path):
        """Convert PDF to HTML"""
        try:
            doc = fitz.open(input_path)
            html_content = """<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>PDF Conversion</title>
    <style>
        body { font-family: Arial, sans-serif; margin: 20px; }
        .page { margin-bottom: 40px; padding: 20px; border: 1px solid #ccc; }
        .page-number { font-weight: bold; color: #666; }
        pre { white-space: pre-wrap; font-family: inherit; }
    </style>
</head>
<body>
    <h1>PDF to HTML Conversion</h1>
"""
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text = page.get_text()
                html_content += f'<div class="page">\n'
                html_content += f'<div class="page-number">Page {page_num + 1}</div>\n'
                html_content += f'<pre>{html.escape(text)}</pre>\n'
                html_content += '</div>\n'
            
            html_content += "</body>\n</html>"
            doc.close()
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(html_content)
            
            return output_path
        except Exception as e:
            raise Exception(f"PDF to HTML conversion failed: {str(e)}")
    
    async def convert_pdf_to_epub(self, input_path, output_path):
        """Convert PDF to EPUB"""
        try:
            # Create a simple EPUB from PDF text
            doc = fitz.open(input_path)
            book = epub.EpubBook()
            
            # Set metadata
            book.set_identifier('pdf_conversion')
            book.set_title('PDF Conversion')
            book.set_language('en')
            
            chapters = []
            
            for page_num in range(min(len(doc), 50)):  # Limit to 50 pages
                page = doc.load_page(page_num)
                text = page.get_text()
                
                # Create chapter
                chapter = epub.EpubHtml(
                    title=f'Page {page_num + 1}',
                    file_name=f'chap_{page_num + 1}.xhtml',
                    lang='en'
                )
                
                chapter.content = f'''
                    <html>
                    <head>
                        <title>Page {page_num + 1}</title>
                    </head>
                    <body>
                        <h1>Page {page_num + 1}</h1>
                        <pre>{html.escape(text)}</pre>
                    </body>
                    </html>
                '''
                
                book.add_item(chapter)
                chapters.append(chapter)
            
            doc.close()
            
            # Define table of contents
            book.toc = chapters
            
            # Add navigation files
            book.add_item(epub.EpubNcx())
            book.add_item(epub.EpubNav())
            
            # Define spine
            book.spine = ['nav'] + chapters
            
            # Write the EPUB file
            epub.write_epub(output_path, book)
            
            return output_path
        except Exception as e:
            raise Exception(f"PDF to EPUB conversion failed: {str(e)}")
    
    async def convert_docx_to_pdf(self, input_path, output_path):
        """Convert DOCX to PDF"""
        try:
            from docx2pdf import convert
            
            convert(input_path, output_path)
            return output_path
        except Exception as e:
            # Fallback using python-docx and reportlab
            return await self._fallback_docx_to_pdf(input_path, output_path)
    
    async def convert_docx_to_txt(self, input_path, output_path):
        """Convert DOCX to TXT"""
        try:
            from docx import Document
            
            doc = Document(input_path)
            text_content = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write('\n'.join(text_content))
            
            return output_path
        except Exception as e:
            raise Exception(f"DOCX to TXT conversion failed: {str(e)}")
    
    async def convert_txt_to_pdf(self, input_path, output_path):
        """Convert TXT to PDF"""
        try:
            with open(input_path, 'r', encoding='utf-8', errors='ignore') as f:
                text_content = f.read()
            
            c = canvas.Canvas(output_path, pagesize=A4)
            width, height = A4
            
            c.setFont("Helvetica", 12)
            y_position = height - 50
            line_height = 14
            
            # Split text into lines
            lines = []
            for paragraph in text_content.split('\n'):
                words = paragraph.split()
                current_line = []
                
                for word in words:
                    test_line = ' '.join(current_line + [word])
                    if c.stringWidth(test_line, "Helvetica", 12) < (width - 100):
                        current_line.append(word)
                    else:
                        if current_line:
                            lines.append(' '.join(current_line))
                        current_line = [word]
                
                if current_line:
                    lines.append(' '.join(current_line))
                lines.append('')
            
            # Add lines to PDF
            for line in lines[:200]:
                if y_position < 50:
                    c.showPage()
                    y_position = height - 50
                    c.setFont("Helvetica", 12)
                
                if line.strip():
                    c.drawString(50, y_position, line)
                y_position -= line_height
            
            c.save()
            return output_path
        except Exception as e:
            raise Exception(f"TXT to PDF conversion failed: {str(e)}")
    
    async def convert_excel_to_pdf(self, input_path, output_path):
        """Convert Excel to PDF"""
        try:
            if input_path.endswith('.csv'):
                df = pd.read_csv(input_path)
            else:
                df = pd.read_excel(input_path)
            
            c = canvas.Canvas(output_path, pagesize=A4)
            width, height = A4
            
            c.setFont("Helvetica-Bold", 16)
            c.drawString(50, height - 50, f"Excel File: {os.path.basename(input_path)}")
            
            c.setFont("Helvetica", 10)
            y_position = height - 80
            
            # Add headers
            headers = df.columns.tolist()
            for i, header in enumerate(headers):
                c.drawString(50 + i * 80, y_position, str(header)[:10])
            
            y_position -= 20
            
            # Add data
            for _, row in df.head(30).iterrows():
                for i, value in enumerate(row):
                    c.drawString(50 + i * 80, y_position, str(value)[:10])
                y_position -= 15
                if y_position < 50:
                    c.showPage()
                    y_position = height - 50
            
            c.save()
            return output_path
        except Exception as e:
            raise Exception(f"Excel to PDF conversion failed: {str(e)}")
    
    async def convert_ppt_to_pdf(self, input_path, output_path):
        """Convert PowerPoint to PDF"""
        try:
            from pptx import Presentation
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import A4
            
            prs = Presentation(input_path)
            c = canvas.Canvas(output_path, pagesize=A4)
            width, height = A4
            
            for i, slide in enumerate(prs.slides):
                if i > 0:
                    c.showPage()
                
                c.setFont("Helvetica-Bold", 16)
                c.drawString(50, height - 50, f"Slide {i+1}")
                
                y_position = height - 80
                c.setFont("Helvetica", 12)
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text[:100]  # Limit text length
                        c.drawString(50, y_position, text)
                        y_position -= 20
                        if y_position < 50:
                            c.showPage()
                            y_position = height - 50
            
            c.save()
            return output_path
        except Exception as e:
            raise Exception(f"PowerPoint to PDF conversion failed: {str(e)}")
    
    async def convert_html_to_pdf(self, input_path, output_path):
        """Convert HTML to PDF"""
        try:
            with open(input_path, 'r', encoding='utf-8', errors='ignore') as f:
                html_content = f.read()
            
            # Simple HTML to text conversion
            import re
            text_content = re.sub('<[^<]+?>', '', html_content)
            
            c = canvas.Canvas(output_path, pagesize=A4)
            width, height = A4
            
            c.setFont("Helvetica", 12)
            y_position = height - 50
            
            lines = text_content.split('\n')
            for line in lines[:100]:
                if y_position < 50:
                    c.showPage()
                    y_position = height - 50
                
                if line.strip():
                    c.drawString(50, y_position, line[:80])
                    y_position -= 15
            
            c.save()
            return output_path
        except Exception as e:
            raise Exception(f"HTML to PDF conversion failed: {str(e)}")
    
    async def convert_markdown_to_pdf(self, input_path, output_path):
        """Convert Markdown to PDF"""
        try:
            with open(input_path, 'r', encoding='utf-8') as f:
                markdown_content = f.read()
            
            # Convert markdown to simple text
            text_content = markdown_content
            
            c = canvas.Canvas(output_path, pagesize=A4)
            width, height = A4
            
            c.setFont("Helvetica", 12)
            y_position = height - 50
            
            lines = text_content.split('\n')
            for line in lines[:100]:
                if y_position < 50:
                    c.showPage()
                    y_position = height - 50
                
                if line.strip():
                    c.drawString(50, y_position, line[:80])
                    y_position -= 15
            
            c.save()
            return output_path
        except Exception as e:
            raise Exception(f"Markdown to PDF conversion failed: {str(e)}")
    
    async def convert_images_to_pdf(self, image_paths, output_path):
        """Convert images to PDF"""
        try:
            valid_images = []
            for img_path in image_paths:
                if os.path.exists(img_path):
                    with Image.open(img_path) as img:
                        img = img.convert('RGB')
                        valid_images.append(img_path)
            
            if not valid_images:
                raise Exception("No valid images found")
            
            with open(output_path, "wb") as f:
                f.write(img2pdf.convert(valid_images))
            
            return output_path
        except Exception as e:
            raise Exception(f"Images to PDF conversion failed: {str(e)}")
    
    async def _fallback_pdf_to_docx(self, input_path, output_path):
        """Fallback PDF to DOCX"""
        try:
            doc = fitz.open(input_path)
            
            from docx import Document
            
            document = Document()
            document.add_heading('PDF Conversion', 0)
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text = page.get_text()
                
                if text.strip():
                    document.add_heading(f'Page {page_num + 1}', level=2)
                    for paragraph in text.split('\n'):
                        if paragraph.strip():
                            document.add_paragraph(paragraph)
            
            doc.close()
            document.save(output_path)
            return output_path
        except Exception as e:
            raise Exception(f"Fallback PDF to DOCX failed: {str(e)}")
    
    async def _fallback_docx_to_pdf(self, input_path, output_path):
        """Fallback DOCX to PDF"""
        try:
            from docx import Document
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import A4
            
            doc = Document(input_path)
            c = canvas.Canvas(output_path, pagesize=A4)
            width, height = A4
            
            c.setFont("Helvetica", 12)
            y_position = height - 50
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    if y_position < 50:
                        c.showPage()
                        y_position = height - 50
                    
                    c.drawString(50, y_position, paragraph.text)
                    y_position -= 15
            
            c.save()
            return output_path
        except Exception as e:
            raise Exception(f"Fallback DOCX to PDF failed: {str(e)}")

# Global converter instance
doc_converter = DocumentConverter()